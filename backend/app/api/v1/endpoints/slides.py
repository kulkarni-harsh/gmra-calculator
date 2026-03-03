import logging
import traceback
from functools import reduce

import pandas as pd
from fastapi import APIRouter, Request
from fastapi.responses import JSONResponse
from geopy.distance import geodesic
from pptx import Presentation

from app.core.config import settings
from app.schemas.slides import GenerateSlidesRequest
from app.services.census import combine_demographics, get_zip_demographics
from app.services.cpt import (
    flag_anchor_cpt_codes,
    generate_cpt_placeholders,
    generate_hospitals_placeholders,
    get_top_cpt_df,
)
from app.services.geocoding import calculate_distance_miles, geocode_addresses, get_location_coordinates
from app.services.mapper import generate_map
from app.services.plots import get_population_distribution_bytes
from app.services.ppt import remove_specific_text_row, replace_all_placeholders, replace_image
from app.services.screenshotter import capture_screen, get_driver
from app.services.specialty import get_specialty_anchor_cpt_info, get_specialty_population
from app.utils.common import get_anchor_cpt_severity_scoring, get_population_severity_scoring

router = APIRouter()


@router.post(
    "/generate",
)
def generate_slides(payload: GenerateSlidesRequest, request: Request):
    filename_prefix = f"outputs/{payload.address_line_1}_{payload.city}_{payload.state}_{payload.zip_code}".replace(
        " ", "_"
    ).replace(".", "_")
    # Get coordinates for input location
    try:
        input_latitude, input_longitude = get_location_coordinates(
            geocoder_client=request.app.state.geocoder_client,
            address_line_1=payload.address_line_1,
            city=payload.city,
            state=payload.state,
            zip_code=payload.zip_code,
        )
    except Exception as e:
        logging.log(logging.ERROR, f"Error geocoding input address: {e}")
        return JSONResponse(content={"message": "Error geocoding input address"}, status_code=400)

    # TODO: Replace with API request to get nearby providers based on input location and specialty
    nearby_providers_df = pd.read_excel("/Users/kulkarni-harsh/Downloads/Frisco OB_GYN.xlsx")

    # Get ZIP code centroids DataFrame from app state
    zip_centroids_df = request.app.state.zip_centroids_df.copy()

    # Geocode nearby provider addresses to get their latitudes and longitudes
    nearby_providers_df = geocode_addresses(nearby_providers_df, request.app.state.geocoder_client)

    # Calculate distances from input location to each provider
    nearby_providers_df["distance_from_source_miles"] = nearby_providers_df.progress_apply(
        lambda r: calculate_distance_miles(
            lat1=float(r["latitude"]),
            lon1=float(r["longitude"]),
            lat2=input_latitude,
            lon2=input_longitude,
        ),
        axis=1,
    )

    # Generate map with nearby providers and save as HTML
    generate_map(
        nearby_providers_df,
        input_latitude,
        input_longitude,
        payload.address_line_1,
        payload.city,
        payload.state,
        payload.zip_code,
        circle_radius_miles=payload.miles_radius,
        html_filepath=filename_prefix + "_map.html",
    )

    # Filter providers within the specified radius and sort by distance
    hospitals_within_range_df = nearby_providers_df[
        nearby_providers_df["distance_from_source_miles"] <= payload.miles_radius
    ].sort_values(by="distance_from_source_miles", ascending=True)
    
    # Log the number of providers found within the radius
    logging.log(
        logging.INFO, f"Found {hospitals_within_range_df.shape[0]} providers within {payload.miles_radius} miles."
    )

    # Take screenshot of the map
    browser = get_driver(width=1600, height=900, scale=2.5)
    try:
        # Use many times without re-installing or re-opening Chrome
        map_bytes = capture_screen(
            browser,
            filename_prefix + "_map.html",
        )
    except Exception as e:
        logging.log(logging.ERROR, f"Error capturing map screenshot: {e}")
        map_bytes = None
    finally:
        # Close ONCE at the very end
        browser.quit()

    # Calculate distances from input location to each ZIP code centroid
    zip_centroids_df["distance_from_source_miles"] = zip_centroids_df.progress_apply(
        lambda row: geodesic((input_latitude, input_longitude), (row["lat"], row["lon"])).miles,
        axis=1,
    )
    # Filter ZIP codes within the specified radius
    filtered_zips_df = zip_centroids_df[zip_centroids_df["distance_from_source_miles"] <= payload.miles_radius]

    # Fallback, if no ZIPs found, take the entered ZIP code
    if filtered_zips_df.empty:
        print("No ZIP codes found within the specified radius. Using the input ZIP code only.")
        filtered_zips_df = pd.DataFrame(
            {
                "zip": [str(payload.zip_code)],
                "lat": [input_latitude],
                "lon": [input_longitude],
                "distance_from_source_miles": [0.0],
            }
        )

    # Get demographic data for the filtered ZIP codes
    zip_demographic_dict = get_zip_demographics(
        tuple(filtered_zips_df["zip"].astype(str).values), settings.CENSUS_API_KEY
    )

    # Combine demographics across all ZIP codes
    combined_demographics_dict = reduce(combine_demographics, zip_demographic_dict.values())

    # Plot population distribution
    population_distribution_bytes = get_population_distribution_bytes(combined_demographics_dict)

    (
        target_demographic_type,
        relevant_population_count,
        target_avg_provider_per_100k,
        actual_avg_provider_per_100k,
        is_provider_ratio_good,
    ) = get_specialty_population(
        hospitals_within_range_df,
        request.app.state.specialty_master_df,
        specialty_name=payload.specialty_name,
        locality_demographics_dict=combined_demographics_dict,
    )

    # Calculate Population Gap analysis verdict
    population_severity_scoring, population_severity_rationale = get_population_severity_scoring(
        current_avg_provider_per_100k=actual_avg_provider_per_100k,
        target_avg_provider_per_100k=target_avg_provider_per_100k,
    )

    # Calculate Clinical Capacity Gap analysis
    (
        (specialty_anchor_cpt_ranges, specialty_anchor_individual_cpt_list),
        target_anchor_cpt_count_per_provider,
    ) = get_specialty_anchor_cpt_info(
        specialty_master_df=request.app.state.specialty_master_df, specialty_name=payload.specialty_name
    )
    # Calculate target anchor visits count
    target_anchor_visits_count = target_anchor_cpt_count_per_provider * hospitals_within_range_df.shape[0]

    # Flag anchor CPT codes and get actual anchor visits count
    hospitals_within_range_df, actual_anchor_visits_count = flag_anchor_cpt_codes(
        hospitals_within_range_df=hospitals_within_range_df,
        specialty_anchor_cpt_ranges=specialty_anchor_cpt_ranges,
        specialty_anchor_individual_cpt_list=specialty_anchor_individual_cpt_list,
    )

    # Calculate Anchor CPT Gap analysis verdict
    (
        anchor_cpt_target_actual_count_diff,
        anchor_cpt_severity_scoring,
        anchor_cpt_severity_rationale,
    ) = get_anchor_cpt_severity_scoring(
        target_anchor_visits_count=target_anchor_visits_count,
        actual_anchor_visits_count=actual_anchor_visits_count,
    )

    # Get Top CPT Codes
    top_cpt_df = get_top_cpt_df(
        hospitals_within_range_df=hospitals_within_range_df, cpt_lookup_df=request.app.state.cpt_lookup_df
    )

    # Create placeholder dict for PPTX
    cpt_placeholder_dict = generate_cpt_placeholders(top_cpt_df)
    hospitals_placeholder_dict, nearest_hospitals_count = generate_hospitals_placeholders(hospitals_within_range_df)

    placeholders_dict = {
        r"{input_address_line_1}": payload.address_line_1,
        r"{input_address_line_2}": payload.address_line_2,
        r"{input_city}": payload.city,
        r"{input_state}": payload.state,
        r"{input_zip_code}": payload.zip_code,
        r"{zip_count}": f"{filtered_zips_df.shape[0]:,}",
        r"{zip_codes_str}": ", ".join(filtered_zips_df["zip"].astype(str).values),
        r"{total_population}": f"{combined_demographics_dict['Total']:,}",
        r"{input_target_specialty}": payload.specialty_name,
        r"{radius}": payload.miles_radius,
        r"{num_providers}": f"{hospitals_within_range_df.shape[0]:,}",
        r"{demographic_type}": target_demographic_type,
        r"{relevant_population_count}": f"{relevant_population_count:,}",
        r"{population_benchmark_density}": f"{target_avg_provider_per_100k:,}",
        r"{population_actual_density}": f"{actual_avg_provider_per_100k:,.1f}",
        r"{population_RAG_status}": population_severity_scoring,
        r"{population_RAG_rationale}": population_severity_rationale,
        # Anchor CPT placeholders
        r"{benchmark_visit_count_per_provider}": f"{target_anchor_cpt_count_per_provider:,}",
        r"{benchmark_visit_total_count}": f"{target_anchor_visits_count:,}",
        r"{actual_visit_total_count}": f"{actual_anchor_visits_count:,}",
        r"{benchmark_actual_visit_total_count_diff}": f"{anchor_cpt_target_actual_count_diff:,}",
        r"{visit_RAG_status}": anchor_cpt_severity_scoring,
        r"{visit_RAG_rationale}": anchor_cpt_severity_rationale,
        **cpt_placeholder_dict,
        **hospitals_placeholder_dict,
    }

    prs = Presentation(settings.TEMPLATES_DIR / "market_gap_analysis_master_slide.pptx")

    for slide in prs.slides:
        # Execute updates
        replace_all_placeholders(slide, placeholders_dict)

    # Replace map on the second slide
    if map_bytes:
        replace_image(prs.slides[1], map_bytes)
        delete_map_slide = False
    else:
        logging.log(logging.WARNING, "Map screenshot not available, skipping map replacement on the slide.")
        delete_map_slide = True

    # Replace population distribution plot on the third slide
    replace_image(prs.slides[2], population_distribution_bytes)

    # Remove rows with 0 CPT counts from the fourth slide
    remove_specific_text_row(
        prs.slides[4],
        col_num=2,
        match_text_list=tuple(["0", 0] + [f"{{CPT_{i}_count}}" for i in range(5, top_cpt_df.shape[0], -1)]),
    )

    # Remove rows with no provider name from the fifth slide
    remove_specific_text_row(
        prs.slides[5],
        col_num=0,
        match_text_list=tuple(f"{{sr_{_}}}" for _ in range(nearest_hospitals_count + 1, 11)),
    )

    if delete_map_slide:
        # If map screenshot is not available, remove the entire slide to avoid empty placeholder
        slide_id_for_deletion = prs.slides[1].slide_id
        slides = prs.slides._sldIdLst
        for sld in slides:
            if sld.id == slide_id_for_deletion:
                slides.remove(sld)
                break
    try:
        prs.save(filename_prefix + ".pptx")
    except Exception:
        traceback.print_exc()
        raise

    hospitals_within_range_df.rename(columns={"distance_from_source_miles": "Distance (miles)"}).to_excel(
        filename_prefix + ".xlsx",
        index=False,
    )

    return JSONResponse(content={"message": "Slides generated successfully"})
