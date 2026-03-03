from pptx.enum.shapes import MSO_SHAPE_TYPE


def replace_all_placeholders(slide, data_map):
    """
    Iterates through all shapes, including tables, to replace placeholders.
    data_map: {"{placeholder}": "replacement_value"}
    """
    for shape in slide.shapes:
        # Standard Text Boxes (Address, Radius, etc.) [cite: 2, 4]
        if shape.has_text_frame:
            _replace_text_in_frame(shape.text_frame, data_map)

        # Table Cells (CPT Codes, Counts, Benchmarks) [cite: 20, 32-46]
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    _replace_text_in_frame(cell.text_frame, data_map)


def _replace_text_in_frame(text_frame, data_map):
    """
    Safely replaces placeholders like {input_city} even if PPT
    has split them into multiple runs. [cite: 2, 12, 15]
    """
    for paragraph in text_frame.paragraphs:
        # 1. Get the full text of the paragraph
        full_text = "".join(run.text for run in paragraph.runs)

        # 2. Check if any of our specific placeholders exist in this paragraph
        updated_text = full_text
        changed = False
        for placeholder, replacement in data_map.items():
            if placeholder in updated_text:
                updated_text = updated_text.replace(placeholder, str(replacement))
                changed = True

        # 3. If a replacement happened, clear runs and reset text
        if changed:
            # Save the formatting of the first run if it exists
            if len(paragraph.runs) > 0:
                # Keep formatting from the first run, clear the rest
                p_run = paragraph.runs[0]
                p_run.text = updated_text
                # Delete subsequent runs to avoid ghost text
                for i in range(1, len(paragraph.runs)):
                    paragraph.runs[i].text = ""
            else:
                paragraph.text = updated_text


def replace_master_map(slide, image_path):
    """
    Identifies the existing placeholder image and replaces it.
    """
    for shape in slide.shapes:
        # Usually, your map will be the largest picture shape or named 'Picture'
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height

            # Remove existing placeholder
            old_img = shape._element
            old_img.getparent().remove(old_img)

            # Add the new analysis map
            slide.shapes.add_picture(image_path, left, top, width, height)
            break


def remove_specific_text_row(slide, col_num, match_text_list: tuple):
    """
    Removes rows from all tables in the slide given a column value.
    slide: pptx slide object
    col_num: zero-based index of the column to check
    match_text_list: list of text values to match for removal (e.g., ["N/A", ""])
    """
    if not match_text_list:
        return
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            # Start from the bottom to avoid index shifting
            for row_idx in range(len(table.rows) - 1, 0, -1):  # skip header row (0)
                cell = table.cell(row_idx, col_num)
                try:
                    value = cell.text.strip()
                    if value in match_text_list:
                        table._tbl.remove(table.rows[row_idx]._tr)
                except Exception as e:
                    print("exception", e)
                    continue
